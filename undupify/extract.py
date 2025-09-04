import io
import os
import zipfile
from typing import List, Tuple

import pandas as pd
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from pypdf import PdfReader


def extract_text_from_txt(path: str) -> str:
	with open(path, 'r', encoding='utf-8', errors='ignore') as f:
		return f.read()


def extract_text_from_pdf(path: str) -> str:
	reader = PdfReader(path)
	parts: List[str] = []
	for page in reader.pages:
		text = page.extract_text() or ''
		parts.append(text)
	return '\n'.join(parts)


def extract_text_from_docx(path: str) -> str:
	doc = Document(path)
	return '\n'.join(p.text for p in doc.paragraphs)


def extract_text_from_pptx(path: str) -> str:
	ppt = Presentation(path)
	texts: List[str] = []
	for slide in ppt.slides:
		for shape in slide.shapes:
			if hasattr(shape, 'text'):
				texts.append(shape.text)
	return '\n'.join(texts)


def extract_text_from_xlsx(path: str) -> str:
	wb = load_workbook(path, read_only=True, data_only=True)
	parts: List[str] = []
	for ws in wb.worksheets:
		for row in ws.iter_rows(values_only=True):
			cells = [str(c) for c in row if c is not None]
			if cells:
				parts.append(' '.join(cells))
	return '\n'.join(parts)


def extract_text(path: str) -> str:
	ext = os.path.splitext(path)[1].lower()
	if ext in {'.txt'}:
		return extract_text_from_txt(path)
	if ext in {'.pdf'}:
		return extract_text_from_pdf(path)
	if ext in {'.docx'}:
		return extract_text_from_docx(path)
	if ext in {'.pptx'}:
		return extract_text_from_pptx(path)
	if ext in {'.xlsx', '.xlsm'}:
		return extract_text_from_xlsx(path)
	# CSV/JSON fallback to DataFrame
	if ext == '.csv':
		df = pd.read_csv(path)
		return '\n'.join(df.astype(str).fillna('').agg(' '.join, axis=1).tolist())
	if ext in {'.json', '.jsonl'}:
		try:
			df = pd.read_json(path, lines=True)
		except ValueError:
			df = pd.read_json(path)
		return '\n'.join(df.astype(str).fillna('').agg(' '.join, axis=1).tolist())
	# Unknown types -> empty
	return ''


def extract_from_zip(zip_path: str) -> List[Tuple[str, str]]:
	items: List[Tuple[str, str]] = []
	with zipfile.ZipFile(zip_path, 'r') as z:
		for info in z.infolist():
			if info.is_dir():
				continue
			name = info.filename
			ext = os.path.splitext(name)[1].lower()
			with z.open(info) as f:
				data = f.read()
				# Store to temp in-memory then parse by type
				buf = io.BytesIO(data)
				tmp_path = name
				# For simplicity: write to disk under a temp folder might be preferable; here, basic support
				# For now, skip complex types inside zip to keep it robust
				text = data.decode('utf-8', errors='ignore') if ext in {'.txt', '.csv', '.json', '.jsonl'} else ''
				items.append((name, text))
	return items


